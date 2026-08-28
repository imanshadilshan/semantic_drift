"""Builds the Day 30 defense/presentation slide deck (PPTX) from the committed IEEE paper
content and the already-generated result figures. No new data, no new claims: every number
here traces to Final Paper/sections/*.tex and Implementation/results/*.csv.

Run from Final Paper/: python scripts/make_slides.py
Writes semantic_drift_defense_slides.pptx.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
FIG_DIR = ROOT / "figures"
OUT_PATH = ROOT / "semantic_drift_defense_slides.pptx"

# Palette, matching the paper's figures.
BLUE = RGBColor(0x1F, 0x3B, 0x5C)      # dark slide background / headers
ACCENT_BLUE = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_ORANGE = RGBColor(0xEB, 0x68, 0x34)
ACCENT_AQUA = RGBColor(0x1B, 0xAF, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF5, 0xF7)
DARK_TEXT = RGBColor(0x22, 0x27, 0x2E)
MID_GRAY = RGBColor(0x5B, 0x63, 0x6E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = color
        shp.line.width = Pt(0.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, italic=False, font=FONT, anchor=None, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=17, color=DARK_TEXT, bullet_color=ACCENT_BLUE,
                 space_after=10, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level, bold = item
        else:
            text, level, bold = item, 0, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.level = 0
        indent = "     " * level
        marker = "-" if level == 0 else u"•"
        run = p.add_run()
        run.text = f"{indent}{marker}  {text}"
        run.font.size = Pt(size - level * 1.5)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color if level == 0 else MID_GRAY
    return tb


def add_footer(slide, page_no, total, section=""):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8), Inches(0.3),
              "Detecting and Mitigating Semantic Drift in Multi-Turn Instruction-Based Image Editing",
              size=9, color=MID_GRAY)
    add_text(slide, Inches(12.3), Inches(7.12), Inches(0.7), Inches(0.3),
              f"{page_no} / {total}", size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def header_bar(slide, kicker, title, title_color=WHITE, kicker_color=ACCENT_ORANGE):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), BLUE)
    add_text(slide, Inches(0.55), Inches(0.14), Inches(12), Inches(0.3), kicker,
              size=13, bold=True, color=kicker_color)
    add_text(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.85), title,
              size=28, bold=True, color=title_color)


def picture_fit(slide, path, x, y, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = Emu(int(max_w / ar))
    else:
        h = max_h
        w = Emu(int(max_h * ar))
    px = Emu(int(x + (max_w - w) / 2))
    py = Emu(int(y + (max_h - h) / 2))
    slide.shapes.add_picture(str(path), px, py, width=w, height=h)


TOTAL_SLIDES = 17

# ------------------------------------------------------------------
# 1. Title slide
# ------------------------------------------------------------------
s = add_slide()
set_background(s, BLUE)
add_rect(s, 0, Inches(4.7), SLIDE_W, Inches(0.06), ACCENT_ORANGE)
add_text(s, Inches(1), Inches(2.15), Inches(11.3), Inches(2.0),
          "Detecting and Mitigating Semantic Drift\nin Multi-Turn Instruction-Based Image Editing",
          size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT, line_spacing=1.05)
add_text(s, Inches(1), Inches(4.95), Inches(10), Inches(0.5),
          "R. W. V. Imansha Dilshan", size=20, bold=True, color=WHITE)
add_text(s, Inches(1), Inches(5.4), Inches(10), Inches(0.8),
          "Software Engineering Teaching Unit\nUniversity of Kelaniya",
          size=15, color=RGBColor(0xC9, 0xD6, 0xE6), line_spacing=1.2)
add_text(s, Inches(1), Inches(6.6), Inches(10), Inches(0.4),
          "Undergraduate Research Project - Final Defense", size=13, italic=True,
          color=ACCENT_ORANGE)

# ------------------------------------------------------------------
# 2. The problem
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "MOTIVATION", "Editors follow the instruction. What else do they change?")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(3.6), [
    "Instruction-based editors (InstructPix2Pix) let users refine an image turn by turn:",
    ("\"make the sky sunset\" then \"add a bird\" then \"remove the fence\"", 1, False),
    "A model can follow every single instruction correctly, and still drift further and",
    "further from what the user wanted, because nothing constrains it to leave the rest",
    "of the image alone.",
    "Existing evaluation (CLIP scores, human preference studies) checks whether the",
    "requested change happened. None of it checks what else changed as a side effect.",
], size=19, space_after=14)
add_rect(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.15), LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(6.02), Inches(11.5), Inches(0.8),
          "This project calls that side effect semantic drift, and asks whether it can be",
          size=18, bold=True, color=BLUE)
add_text(s, Inches(0.9), Inches(6.42), Inches(11.5), Inches(0.5),
          "measured automatically and reduced cheaply.", size=18, bold=True, color=BLUE)
add_footer(s, 2, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 3. Research questions
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "SCOPE", "Three research questions")
rq_data = [
    ("RQ1", ACCENT_BLUE, "Can unintended drift be measured automatically,",
     "without manually inspecting every edited image?"),
    ("RQ2", ACCENT_ORANGE, "Does drift accumulate as an edit chain gets longer,",
     "so later instructions cause more collateral damage than earlier ones?"),
    ("RQ3", ACCENT_AQUA, "Can cheap, pretrained-model-only interventions",
     "reduce that drift?"),
]
y = Inches(1.75)
for tag, color, l1, l2 in rq_data:
    add_rect(s, Inches(0.6), y, Inches(1.7), Inches(1.45), color)
    add_text(s, Inches(0.6), y + Inches(0.42), Inches(1.7), Inches(0.6), tag, size=26,
              bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.5), y + Inches(0.08), Inches(10.2), Inches(0.6), l1, size=19, bold=True)
    add_text(s, Inches(2.5), y + Inches(0.62), Inches(10.2), Inches(0.6), l2, size=17,
              color=MID_GRAY)
    y += Inches(1.7)
add_text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.4),
          "Scoped for a single contributor, 30 days, pretrained models only, free-tier compute.",
          size=14, italic=True, color=MID_GRAY)
add_footer(s, 3, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 4. Gap in prior work
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "RELATED WORK", "What prior approaches do not measure")
rows = [
    ("Prompt-to-Prompt / Plug-and-Play / MasaCtrl / Null-text Inv.", "Produce a more consistent single edit", False),
    ("Emu Edit / HIVE", "Improve whether the requested edit happened", False),
    ("MagicBrush", "Multi-turn data, but only checks the target region changed correctly", False),
    ("This work (Drift Score)", "Measures unintended change, across a chain, with mitigations compared", True),
]
y = Inches(1.85)
for name, desc, highlight in rows:
    bg = ACCENT_BLUE if highlight else LIGHT_GRAY
    tcol = WHITE if highlight else DARK_TEXT
    dcol = RGBColor(0xDD, 0xE8, 0xF7) if highlight else MID_GRAY
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.05), bg)
    add_text(s, Inches(0.85), y + Inches(0.08), Inches(11.6), Inches(0.42), name, size=17,
              bold=True, color=tcol)
    add_text(s, Inches(0.85), y + Inches(0.52), Inches(11.6), Inches(0.42), desc, size=14,
              color=dcol)
    y += Inches(1.2)
add_footer(s, 4, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 5. The Drift Score
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "METHODOLOGY", "The Drift Score")
picture_fit(s, FIG_DIR / "fig1_pipeline.png", Inches(0.6), Inches(1.6), Inches(12.1), Inches(2.1))
add_bullets(s, Inches(0.6), Inches(3.85), Inches(12.1), Inches(2.9), [
    "Segment the pre-edit image with SAM into regions.",
    "Identify the target region: the box CLIP finds most similar to the instruction text.",
    "For every OTHER (non-target) region, measure 1 minus CLIP cosine similarity between",
    ("its pre-edit and post-edit crop.", 1, False),
    "Per-step Drift Score = mean drift over non-target regions.",
    "Cumulative Drift Score = sum of per-step scores across a chain. Higher = more",
    ("unintended change accumulated.", 1, False),
], size=17, space_after=8)
add_footer(s, 5, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 6. Dataset
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "METHODOLOGY", "Dataset: 120 four-step edit chains")
stat_data = [
    ("60", "COCO val2017 images", ACCENT_BLUE),
    ("120", "edit chains, 4 steps each", ACCENT_ORANGE),
    ("480", "individual edit instructions", ACCENT_AQUA),
]
x = Inches(0.6)
for num, label, color in stat_data:
    add_rect(s, x, Inches(1.75), Inches(3.9), Inches(1.7), LIGHT_GRAY)
    add_text(s, x, Inches(1.9), Inches(3.9), Inches(0.9), num, size=44, bold=True, color=color,
              align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.85), Inches(3.9), Inches(0.5), label, size=14, color=MID_GRAY,
              align=PP_ALIGN.CENTER)
    x += Inches(4.1)
add_bullets(s, Inches(0.6), Inches(3.75), Inches(12.1), Inches(3.0), [
    "Chain A (object-level): 4 localized, single-object edits per image",
    ("e.g. change a horse's color, then remove an obstacle in front of it", 1, False),
    "Chain B (global): 4 broad, scene-level edits per image",
    ("e.g. make the scene rainy, then add sunset lighting", 1, False),
    "Every image inspected individually and hand-written, not template-generated",
    "Instruction verbs (all 480): change 108, add 114, remove 60, global/stylistic 198",
    ("(hand-written by a single author, a stated bias, not crowd-sourced)", 1, False),
], size=16, space_after=8)
add_footer(s, 6, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 7. Mitigations
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "METHODOLOGY", "Two mitigation strategies")
add_rect(s, Inches(0.6), Inches(1.75), Inches(5.9), Inches(4.6), LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(1.95), Inches(5.3), Inches(0.5), "Region-Locking",
          size=20, bold=True, color=ACCENT_ORANGE)
add_bullets(s, Inches(0.9), Inches(2.55), Inches(5.3), Inches(3.6), [
    "Run the edit on the full image as usual.",
    "Revert every pixel outside the target region's box back to the pre-edit image.",
    "A correction applied AFTER generation.",
], size=16, space_after=12)
add_rect(s, Inches(6.8), Inches(1.75), Inches(5.9), Inches(4.6), LIGHT_GRAY)
add_text(s, Inches(7.1), Inches(1.95), Inches(5.3), Inches(0.5), "Masked Conditioning",
          size=20, bold=True, color=ACCENT_AQUA)
add_bullets(s, Inches(7.1), Inches(2.55), Inches(5.3), Inches(3.6), [
    "Crop the image to the target region, plus a 32-pixel context margin.",
    "Generate the edit only on that crop.",
    "Paste the result back into the full frame. The model never sees the rest of the",
    ("image, a constraint applied AT generation time.", 1, False),
], size=16, space_after=12)
add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
          "Both use the identical SAM + CLIP target-region procedure the Drift Score itself uses.",
          size=14, italic=True, color=MID_GRAY)
add_footer(s, 7, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 8. Edit Adherence Score
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "METHODOLOGY", "A second, independent check: Edit Adherence Score")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.4), [
    "Low unintended change only matters if the requested edit still happened.",
    "Edit Adherence Score = whole-image CLIP similarity between the final image and the",
    ("final instruction's text.", 1, False),
    "Computed identically for all three conditions, with no region segmentation at all,",
    "so it cannot inherit the target-box definition the Drift Score and region-locking share.",
], size=18, space_after=10)
add_rect(s, Inches(0.6), Inches(4.35), Inches(12.1), Inches(1.7), RGBColor(0xFB, 0xEE, 0xE5))
add_text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(0.5),
          "Why this matters: region-locking's low drift score is EXPECTED almost by definition,",
          size=17, bold=True, color=ACCENT_ORANGE)
add_text(s, Inches(0.9), Inches(5.05), Inches(11.5), Inches(0.5),
          "since it reverts pixels using the same box the score excludes. This metric answers",
          size=17, bold=True, color=ACCENT_ORANGE)
add_text(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.4),
          "that concern directly, with independent evidence.", size=17, bold=True, color=ACCENT_ORANGE)
add_footer(s, 8, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 9. RQ1: validating the metric
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "RESULTS - RQ1", "Does the score agree with what the eye sees?")
picture_fit(s, FIG_DIR / "fig5_surfer_example.png", Inches(0.6), Inches(1.6), Inches(6.0), Inches(3.3))
add_text(s, Inches(0.6), Inches(5.0), Inches(6.0), Inches(0.9),
          "Lowest-scoring chain (0.082): only the outfit and\nsurfboard color change; wave, spray, and sky untouched.",
          size=14, color=MID_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)
picture_fit(s, FIG_DIR / "fig3_qualitative_example.png", Inches(6.8), Inches(1.6), Inches(5.9), Inches(3.3))
add_text(s, Inches(6.8), Inches(5.0), Inches(5.9), Inches(0.9),
          "Highest-scoring chain (0.565): a severe editing\nfailure, collapses to near-black by step 3.",
          size=14, color=MID_GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)
add_text(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.9),
          "The metric correctly ranks the best and worst chains in the dataset, agreement between the",
          size=15, bold=True, color=BLUE)
add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.5),
          "number and visible content that a purely statistical validation cannot substitute for.",
          size=15, bold=True, color=BLUE)
add_footer(s, 9, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 10. RQ3: mitigation effectiveness
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "RESULTS - RQ3", "Both mitigations cut drift by a wide margin")
picture_fit(s, FIG_DIR / "fig2_drift_by_condition.png", Inches(0.6), Inches(1.6), Inches(6.4), Inches(4.6))
tbl_x, tbl_y = Inches(7.4), Inches(2.1)
headers = ["Condition", "Mean", "Reduction", "p", "dz"]
data = [
    ["Baseline", "0.293", "-", "-", "-"],
    ["Masked Cond.", "0.092", "68.6%", "<0.0001", "1.73"],
    ["Region-Locking", "0.032", "89.1%", "<0.0001", "2.68"],
]
col_w = [Inches(1.9), Inches(0.95), Inches(1.15), Inches(1.15), Inches(0.75)]
row_h = Inches(0.62)
x = tbl_x
for i, hd in enumerate(headers):
    add_rect(s, x, tbl_y, col_w[i], row_h, BLUE)
    add_text(s, x, tbl_y + Inches(0.13), col_w[i], Inches(0.4), hd, size=13, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    x += col_w[i]
for r, row in enumerate(data):
    y = tbl_y + row_h * (r + 1)
    rowbg = LIGHT_GRAY if r % 2 == 0 else WHITE
    x = tbl_x
    for i, val in enumerate(row):
        add_rect(s, x, y, col_w[i], row_h, rowbg)
        add_text(s, x, y + Inches(0.13), col_w[i], Inches(0.4), val, size=13,
                  bold=(i == 0), align=PP_ALIGN.CENTER)
        x += col_w[i]
add_text(s, Inches(7.4), Inches(4.78), Inches(5.5), Inches(1.6),
          "Region-locking also beats masked conditioning\ndirectly (p<0.0001, dz=1.05). Rank-biserial\ncorrelation = 1.00: every one of 120 paired\nchains favored region-locking.",
          size=15, color=MID_GRAY, line_spacing=1.2)
add_footer(s, 10, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 11. Edit adherence trade-off
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "RESULTS", "Real, but modest, cost to instruction-following")
tbl_x, tbl_y = Inches(0.8), Inches(1.85)
headers = ["Condition", "CLIP Score", "vs. Baseline", "p", "dz"]
data = [
    ["Baseline (pre-edit)", "0.230", "-", "-", "-"],
    ["Baseline", "0.245", "-", "-", "-"],
    ["Masked Conditioning", "0.225", "-0.020", "<0.0001", "-0.51"],
    ["Region-Locking", "0.217", "-0.028", "<0.0001", "-0.68"],
]
col_w = [Inches(3.1), Inches(1.9), Inches(1.9), Inches(1.9), Inches(1.4)]
row_h = Inches(0.62)
x = tbl_x
for i, hd in enumerate(headers):
    add_rect(s, x, tbl_y, col_w[i], row_h, BLUE)
    add_text(s, x, tbl_y + Inches(0.13), col_w[i], Inches(0.4), hd, size=14, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    x += col_w[i]
for r, row in enumerate(data):
    y = tbl_y + row_h * (r + 1)
    rowbg = LIGHT_GRAY if r % 2 == 0 else WHITE
    x = tbl_x
    for i, val in enumerate(row):
        add_rect(s, x, y, col_w[i], row_h, rowbg)
        add_text(s, x, y + Inches(0.13), col_w[i], Inches(0.4), val, size=14,
                  bold=(i == 0), align=PP_ALIGN.CENTER)
        x += col_w[i]
add_rect(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.65), RGBColor(0xEA, 0xF6, 0xF0))
add_bullets(s, Inches(1.1), Inches(5.25), Inches(11.1), Inches(1.5), [
    "Both mitigations do cost real instruction-adherence (p<0.0001, not zero).",
    "But the cost's effect size (dz 0.51 to 0.68) is under half the size of the matching",
    ("drift-reduction effect size (dz 1.73 to 2.68).", 1, False),
    "A genuine, quantified trade-off that favors both mitigations, not a free lunch,",
    ("and not circular reasoning either.", 1, False),
], size=14, color=DARK_TEXT, space_after=5)
add_footer(s, 11, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 12. RQ2: compounding
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "RESULTS - RQ2", "Drift did not compound the way we expected")
picture_fit(s, FIG_DIR / "fig4_step_position.png", Inches(0.6), Inches(1.6), Inches(6.4), Inches(4.5))
add_bullets(s, Inches(7.3), Inches(1.9), Inches(5.5), Inches(4.2), [
    "No single consecutive step transition is significant",
    ("(1→2: p=0.85, 2→3: p=0.33, 3→4: p=0.47)", 1, False),
    "Step 1 vs. step 4 IS significant (p=0.012), but",
    ("step 4 shows LESS drift than step 1, the reverse", 1, False),
    ("of the predicted direction", 1, False),
    "Explanation: a chain that collapses catastrophically",
    ("early has little room left to register as further", 1, False),
    ("changed, a measurement ceiling effect, not proof", 1, False),
    ("that later edits are genuinely gentler", 1, False),
], size=16, space_after=10)
add_footer(s, 12, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 13. Why the hypotheses reversed
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "DISCUSSION", "Two hypotheses reversed. Both have a mechanism.")
add_rect(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(2.3), LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(1.85), Inches(11.5), Inches(0.45),
          "H2: masked conditioning was predicted to beat region-locking. It didn't.",
          size=17, bold=True, color=ACCENT_ORANGE)
add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.5),
          "Masked conditioning edits a padded crop (target box + 32px margin) that the Drift\n"
          "Score does NOT exclude from scoring, so its own padding ring counts as drift even\n"
          "when nothing objectionable happened there. Region-locking's hard revert has no\n"
          "equivalent margin.", size=15, color=DARK_TEXT, line_spacing=1.25)
add_rect(s, Inches(0.6), Inches(4.25), Inches(12.1), Inches(2.3), LIGHT_GRAY)
add_text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.45),
          "H1: drift was predicted to compound across the chain. It didn't, measurably.",
          size=17, bold=True, color=ACCENT_ORANGE)
add_text(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.5),
          "A shared weak point for both mitigations: CLIP target-region identification is\n"
          "unreliable for \"add X\" instructions, since nothing pre-existing matches well. One\n"
          "traced failure: an \"add a cap\" box covered 99.6% of the frame, giving neither\n"
          "mitigation a meaningful constraint on that step.", size=15, color=DARK_TEXT, line_spacing=1.25)
add_footer(s, 13, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 14. Limitations
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "LIMITATIONS", "Seven constraints on the present study")
add_bullets(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(4.8), [
    "60-image, 120-chain dataset: an initial empirical study, not a definitive benchmark",
    "Instructions hand-written by a single author, not crowd-sourced (a stated bias)",
    "All results specific to InstructPix2Pix; model-agnosticism is a design goal, untested",
    "CLIP similarity is an imperfect proxy for human-perceived change",
    "The human-rating tool (15 stratified chains) was built but not administered",
    "Target-region identification is a known weak point for \"add\" instructions",
    "Attention-restricted editing was scoped but not implemented (RQ3 answered without it)",
], size=18, space_after=16)
add_footer(s, 14, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 15. Conclusion
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "CONCLUSION", "RQ1 and RQ3: yes. RQ2: more interesting than yes/no.")
add_bullets(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(4.5), [
    "RQ1: unintended drift CAN be measured automatically, and the metric agrees with",
    ("visible content on both the best and worst chains in the dataset.", 1, False),
    "RQ3: cheap, pretrained-model-only mitigations DO work. Region-locking cuts drift",
    ("89.1%, masked conditioning 68.6%, both p<0.0001, both large effect sizes.", 1, False),
    "RQ2: no reliable evidence of compounding. Drift instead looks front-loaded into",
    ("occasional catastrophic failures, a distinct and equally useful finding.", 1, False),
    "Two specific mechanisms explain both reversed hypotheses, evidence the metric",
    ("captures something real and interpretable, not noise.", 1, False),
], size=18, space_after=12)
add_footer(s, 15, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 16. Future work
# ------------------------------------------------------------------
s = add_slide()
set_background(s, WHITE)
header_bar(s, "FUTURE WORK", "Next steps toward a benchmark")
add_bullets(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(4.5), [
    "Run the built human-perceptual check to validate the Drift Score and the Edit",
    ("Adherence Score against human judgment on the same 15 chains.", 1, False),
    "Evaluate the Drift Score against a second, more recent editor, to test the",
    ("model-agnosticism claimed but not exercised here.", 1, False),
    "Design a metric less prone to the ceiling effect, sensitive to further degradation",
    ("of an already-degraded image, to separate gradual compounding from sudden failure.", 1, False),
], size=19, space_after=16)
add_footer(s, 16, TOTAL_SLIDES)

# ------------------------------------------------------------------
# 17. Thank you
# ------------------------------------------------------------------
s = add_slide()
set_background(s, BLUE)
add_rect(s, 0, Inches(3.55), SLIDE_W, Inches(0.06), ACCENT_ORANGE)
add_text(s, Inches(1), Inches(2.6), Inches(11.3), Inches(0.9), "Thank you",
          size=44, bold=True, color=WHITE)
add_text(s, Inches(1), Inches(3.75), Inches(11.3), Inches(0.5), "Questions?",
          size=22, color=RGBColor(0xC9, 0xD6, 0xE6))
add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
          "Code, dataset, and full paper: github.com/imanshadilshan/semantic_drift",
          size=14, color=ACCENT_ORANGE)

prs.save(OUT_PATH)
print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides)")
